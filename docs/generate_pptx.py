#!/usr/bin/env python3
"""Generate QuizMaster project presentation as PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BLUE_700   = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_500   = RGBColor(0x3B, 0x82, 0xF6)
PURPLE_600 = RGBColor(0x7E, 0x22, 0xCE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_DARK   = RGBColor(0x1E, 0x29, 0x3B)
TEXT_GRAY   = RGBColor(0x64, 0x74, 0x8B)
GREEN_TXT  = RGBColor(0x16, 0x65, 0x34)
RED_TXT    = RGBColor(0xDC, 0x26, 0x26)
GREEN_OK   = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=None, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    if color: p.font.color.rgb = color
    return tf

def add_para(tf, text, font_size=18, color=None, bold=False, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = "Arial"
    p.space_before = Pt(space_before)
    if color: p.font.color.rgb = color
    return p

def add_table(slide, left, top, width, rows_data, col_widths=None, header_bg=None):
    rows, cols = len(rows_data), len(rows_data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.35 * rows))
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths): t.columns[i].width = Inches(w)
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(val)
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13); p.font.name = "Arial"
            if r == 0 and header_bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True; p.font.color.rgb = WHITE; p.font.size = Pt(12)
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return ts

def add_page_number(slide, num, total=9, color=None):
    add_textbox(slide, 12.1, 7.0, 1.0, 0.4, f"{num} / {total}", font_size=11, color=color or TEXT_GRAY, align=PP_ALIGN.RIGHT)

# ═══ SLIDE 1: COVER ═══
s = add_blank_slide()
add_bg(s, RGBColor(0x1E, 0x3A, 0x8A))
add_textbox(s, 0.5, 0.3, 3.0, 0.5, "\U0001F3AF QuizMaster", font_size=16, color=RGBColor(0x99,0x99,0xBB), bold=True)
add_textbox(s, 0.7, 2.0, 10.0, 1.5, "Interactive Real-Time\nQuiz Platform", font_size=52, bold=True, color=WHITE)
add_textbox(s, 0.7, 3.8, 8.0, 0.6, "Web Application MVP — Project Presentation", font_size=24, color=RGBColor(0xBB,0xBB,0xDD))
add_textbox(s, 0.7, 5.0, 5.0, 0.4, "Stack:  React · Express · Socket.IO · SQLite", font_size=16, color=RGBColor(0x99,0x99,0xBB))
add_textbox(s, 0.7, 5.5, 5.0, 0.4, "Development:  ~5 Days · 39 Source Files", font_size=16, color=RGBColor(0x99,0x99,0xBB))
add_page_number(s, 1, color=RGBColor(0x88,0x88,0xAA))

# ═══ SLIDE 2: OVERVIEW ═══
s = add_blank_slide()
add_bg(s, LIGHT_GRAY)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "Project Overview", font_size=36, bold=True, color=TEXT_DARK)
add_textbox(s, 0.7, 1.1, 11.0, 0.5, "A full-stack web application for creating and conducting live quizzes with real-time interaction.", font_size=18, color=TEXT_GRAY)

cards = [
    ("\U0001F4DD", "Quiz Creation", "4 question types: text & image,\nsingle & multiple choice."),
    ("⚡", "Real-Time Play", "Questions sync via WebSocket.\nCountdown timers per question."),
    ("\U0001F3C6", "Live Leaderboard", "Auto-scoring. Final rankings with\nscore, accuracy & response time."),
    ("\U0001F511", "Room Codes", "6-character codes for instant\njoining. No links, no invites."),
    ("\U0001F464", "Dual Roles", "Organizer dashboard + Participant\ninterface for seamless gameplay."),
    ("\U0001F4CA", "History & Stats", "Personal profiles with participation\nhistory & performance metrics."),
]
for i, (icon, title, body) in enumerate(cards):
    left = 0.7 + (i % 3) * 4.1
    top = 2.2 + (i // 3) * 2.6
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.8), Inches(2.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]; p.text = f"{icon}  {title}"
    p.font.size = Pt(16); p.font.bold = True; p.font.name = "Arial"; p.font.color.rgb = TEXT_DARK
    p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(13); p2.font.name = "Arial"
    p2.font.color.rgb = TEXT_GRAY; p2.space_before = Pt(6)
add_page_number(s, 2, color=TEXT_GRAY)

# ═══ SLIDE 3: ARCHITECTURE ═══
s = add_blank_slide()
add_bg(s, BG_DARK)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "System Architecture", font_size=36, bold=True, color=WHITE)
add_textbox(s, 0.7, 1.8, 5.5, 0.5, "\U0001F5A5  Frontend", font_size=22, bold=True, color=BLUE_500)
add_table(s, 0.7, 2.4, 5.5, [
    ["Component","Choice"],["Framework","React 18"],["Build Tool","Vite 5"],
    ["Styling","Tailwind CSS 3"],["Routing","React Router 6"],
    ["Real-Time","Socket.IO Client"],["HTTP","Fetch API"]
], [2.5,3.0], BLUE_700)
add_textbox(s, 7.0, 1.8, 5.5, 0.5, "⚙  Backend", font_size=22, bold=True, color=BLUE_500)
add_table(s, 7.0, 2.4, 5.5, [
    ["Component","Choice"],["Runtime","Node.js"],["Framework","Express 4"],
    ["Database","SQLite (better-sqlite3)"],["Real-Time","Socket.IO 4"],
    ["Auth","JWT (jsonwebtoken)"],["Validation","express-validator"]
], [2.5,3.0], BLUE_700)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.1))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E,0x3A,0x5A); box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.text = "Key Decision:  Single language (JavaScript) across the entire stack. SQLite requires zero configuration — ideal for rapid prototyping."
p.font.size = Pt(15); p.font.name = "Arial"; p.font.color.rgb = RGBColor(0xBB,0xCC,0xDD)
add_page_number(s, 3)

# ═══ SLIDE 4: TECHNOLOGY RATIONALE ═══
s = add_blank_slide()
add_bg(s, LIGHT_GRAY)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "Technology Rationale", font_size=36, bold=True, color=TEXT_DARK)
add_textbox(s, 0.7, 1.7, 5.5, 0.5, "✅  Chosen", font_size=20, bold=True, color=GREEN_OK)
add_table(s, 0.7, 2.3, 5.8, [
    ["Technology","Rationale"],["Express","Lightweight, vast middleware ecosystem"],
    ["SQLite","Zero-config, file-based, sync API"],
    ["Socket.IO","Auto-reconnect, rooms, polling fallback"],
    ["Vite","Instant HMR, 10× faster than Webpack"],
    ["Tailwind CSS","Utility-first, tiny production bundle (~5KB)"],
    ["JWT","Stateless, works with WebSocket auth"]
], [2.2,3.6], GREEN_TXT)
add_textbox(s, 7.0, 1.7, 5.5, 0.5, "❌  Not Chosen", font_size=20, bold=True, color=RED_TXT)
add_table(s, 7.0, 2.3, 5.8, [
    ["Alternative","Reason for Rejection"],["Next.js","Overkill for SPA; adds SSR complexity"],
    ["MongoDB","External server needed; data is relational"],
    ["PostgreSQL","Setup overhead for MVP stage"],
    ["Redux","Excessive boilerplate; Context API suffices"],
    ["Webpack","Slower dev server; complex configuration"],
    ["SASS","Separate files; slower iteration"]
], [2.2,3.6], RED_TXT)
add_page_number(s, 4, color=TEXT_GRAY)

# ═══ SLIDE 5: DATABASE ═══
s = add_blank_slide()
add_bg(s, BG_DARK)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "Database Model", font_size=36, bold=True, color=WHITE)
add_textbox(s, 0.7, 1.7, 5.5, 0.5, "Core Entities (6 Tables)", font_size=20, bold=True, color=BLUE_500)
add_table(s, 0.7, 2.3, 5.5, [
    ["Table","Purpose"],["users","Authentication, role-based access"],
    ["quizzes","Quiz metadata, settings, ownership"],
    ["questions","4 types, points, time limits, ordering"],
    ["options","Answer choices with correctness flag"],
    ["quiz_sessions","Room codes, lifecycle state machine"],
    ["participants","Join records, cumulative scores"],
    ["answers","Response tracking, grading, timing"]
], [2.5,3.0], BLUE_700)
add_textbox(s, 7.0, 1.7, 5.5, 0.5, "Key Design Decisions", font_size=20, bold=True, color=BLUE_500)
tf = add_textbox(s, 7.0, 2.3, 5.5, 4.0, "", font_size=14, color=RGBColor(0xCC,0xCC,0xDD))
decisions = [
    "Foreign keys with CASCADE for clean deletion",
    "JSON column for selected_option_ids",
    "Room codes exclude 0/O and 1/I",
    "WAL mode for concurrent performance",
    "Indexes on all foreign keys"
]
for i, d in enumerate(decisions):
    if i == 0:
        tf.paragraphs[0].text = f"▸  {d}"
        tf.paragraphs[0].font.size = Pt(14); tf.paragraphs[0].font.color.rgb = RGBColor(0xCC,0xCC,0xDD)
    else:
        add_para(tf, f"▸  {d}", font_size=14, color=RGBColor(0xCC,0xCC,0xDD), space_before=8)
add_page_number(s, 5)

# ═══ SLIDE 6: API & WEBSOCKET ═══
s = add_blank_slide()
add_bg(s, LIGHT_GRAY)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "API & WebSocket Design", font_size=36, bold=True, color=TEXT_DARK)
add_textbox(s, 0.7, 1.7, 5.5, 0.5, "REST API — 14 Endpoints", font_size=20, bold=True, color=BLUE_500)
add_table(s, 0.7, 2.3, 5.8, [
    ["Method","Endpoint","Description"],
    ["POST","/api/auth/register","Create account with role"],
    ["POST","/api/auth/login","Get JWT token"],
    ["GET","/api/quizzes","List organizer's quizzes"],
    ["POST","/api/quizzes","Create new quiz"],
    ["PUT","/api/quizzes/:id","Update quiz settings"],
    ["POST","/api/quizzes/:id/questions","Add question + options"],
    ["POST","/api/sessions","Create room (organizer)"],
    ["POST","/api/sessions/:code/join","Join as participant"],
    ["GET","/api/users/me/history","Participation history"],
], [1.0,3.0,1.8], BLUE_700)
add_textbox(s, 7.0, 1.7, 5.5, 0.5, "WebSocket — 11 Events", font_size=20, bold=True, color=BLUE_500)
add_table(s, 7.0, 2.3, 5.8, [
    ["Event","Direction","Purpose"],
    ["session:start","Org → Server","Begin quiz"],
    ["question:show","Server → All","Push question + timer"],
    ["answer:submit","Player → Server","Submit response"],
    ["question:closed","Server → All","Timeout / next"],
    ["quiz:leaderboard","Server → All","Live rankings"],
    ["quiz:finished","Server → All","Final results"],
    ["organizer:stats","Server → Org","Live answer counts"],
    ["answer:result","Server → Player","Per-player feedback"],
], [2.2,1.8,1.8], BLUE_700)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.0))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xEF,0xF6,0xFF); box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.text = "Server-enforced: Late answers rejected. Timer managed server-side. One answer per participant per question. Room-code-scoped socket rooms."
p.font.size = Pt(14); p.font.name = "Arial"; p.font.color.rgb = TEXT_DARK
add_page_number(s, 6, color=TEXT_GRAY)

# ═══ SLIDE 7: DEVELOPMENT STAGES ═══
s = add_blank_slide()
add_bg(s, BG_DARK)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "Development Stages", font_size=36, bold=True, color=WHITE)
stages = [
    ("1", "Design & Spec", "Complete spec-as-code.\nData models, API contracts,\nWebSocket protocol, UI states.", "1 Day"),
    ("2", "Backend", "JWT auth, CRUD APIs,\nsession management, room\ncodes, WebSocket handler.", "1 Day"),
    ("3", "Frontend", "10 pages: auth, dashboard,\nquiz editor, lobby, player\nroom, profile. Real-time UI.", "2 Days"),
    ("4", "Integration", "E2E API testing. Production\nbuild. CORS & WebSocket\nsetup. Documentation.", "1 Day"),
]
for i, (num, title, body, days) in enumerate(stages):
    left = 0.7 + i * 3.15
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(2.9), Inches(4.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x22,0x22,0x3E); shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]; p.text = f"{num}  {title}"
    p.font.size = Pt(20); p.font.bold = True; p.font.name = "Arial"; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(13); p2.font.name = "Arial"
    p2.font.color.rgb = RGBColor(0xAA,0xAA,0xBB); p2.space_before = Pt(10)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left+0.15), Inches(5.6), Inches(1.2), Inches(0.35))
    badge.fill.solid(); badge.fill.fore_color.rgb = BLUE_500; badge.line.fill.background()
    btf = badge.text_frame; btf.paragraphs[0].text = days
    btf.paragraphs[0].font.size = Pt(12); btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].font.color.rgb = WHITE; btf.paragraphs[0].alignment = PP_ALIGN.CENTER
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.55))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E,0x3A,0x5A); box.line.fill.background()
tf = box.text_frame; tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]; p.text = "Parallel-first: Spec drove simultaneous backend & frontend generation. Zero rework — API contracts matched from start."
p.font.size = Pt(13); p.font.name = "Arial"; p.font.color.rgb = RGBColor(0x99,0xAA,0xBB)
add_page_number(s, 7)

# ═══ SLIDE 8: METRICS ═══
s = add_blank_slide()
add_bg(s, PURPLE_600)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "Project Metrics", font_size=36, bold=True, color=WHITE)
stats = [("39","Source Files"),("~4,500","Lines of Code"),("14","API Endpoints"),("11","WS Events"),("10","UI Pages")]
for i, (val, label) in enumerate(stats):
    left = 0.7 + i * 2.45
    add_textbox(s, left, 2.5, 2.2, 1.0, val, font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, left, 3.5, 2.2, 0.4, label, font_size=14, color=RGBColor(0xCC,0xCC,0xEE), align=PP_ALIGN.CENTER)
for left, title, lines in [
    (0.7, "Production Bundle", ["JavaScript:  257 KB  (80 KB gzip)", "CSS:            28 KB  (5 KB gzip)"]),
    (6.8, "Quality", ["✅  0 build errors", "✅  79 modules bundled", "✅  Full API test pass"])
]:
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.5), Inches(5.8), Inches(1.8))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    tf = card.text_frame; tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(18); p.font.bold = True; p.font.name = "Arial"; p.font.color.rgb = TEXT_DARK
    for line in lines: add_para(tf, line, font_size=16, color=TEXT_GRAY, space_before=8)
add_page_number(s, 8, color=RGBColor(0xCC,0xCC,0xEE))

# ═══ SLIDE 9: FEATURES & LINKS ═══
s = add_blank_slide()
add_bg(s, LIGHT_GRAY)
add_textbox(s, 0.7, 0.5, 11.0, 0.7, "Feature Checklist & Links", font_size=36, bold=True, color=TEXT_DARK)
add_textbox(s, 0.7, 1.7, 5.5, 0.5, "Requirements Met", font_size=20, bold=True, color=GREEN_OK)
tf = add_textbox(s, 0.7, 2.3, 5.5, 3.5, "")
for i, item in enumerate([
    "✅  User registration with roles",
    "✅  Quiz CRUD with categories & settings",
    "✅  4 question types (text/image × single/multi)",
    "✅  Room code join system (6-char codes)",
    "✅  Real-time sync via WebSocket",
    "✅  Countdown timer (server-enforced)",
    "✅  Auto-scoring & live leaderboard",
    "✅  Personal history & statistics",
    "✅  Responsive design (Tailwind CSS)",
]):
    if i == 0:
        tf.paragraphs[0].text = item; tf.paragraphs[0].font.size = Pt(15); tf.paragraphs[0].font.name = "Arial"; tf.paragraphs[0].font.color.rgb = TEXT_DARK
    else:
        add_para(tf, item, font_size=15, color=TEXT_DARK, space_before=5)

add_textbox(s, 7.0, 1.7, 5.5, 0.5, "Links", font_size=20, bold=True, color=BLUE_500)
for i, (title, url) in enumerate([
    ("\U0001F4C2  Repository", "[GitHub URL]"),
    ("\U0001F4CB  Specification", "spec.md in repository root"),
    ("\U0001F3A8  Mockups", "[Figma / Miro link]"),
    ("\U0001F680  Live Demo", "[Deployment URL]"),
]):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.3 + i*1.0), Inches(5.5), Inches(0.85))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    tf = card.text_frame; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.name = "Arial"; p.font.color.rgb = TEXT_DARK
    add_para(tf, url, font_size=12, color=TEXT_GRAY, space_before=4)
add_textbox(s, 0.7, 6.5, 11.9, 0.6, "\U0001F3AF  Ready to run: npm run dev  —  Server :3000, Client :5173", font_size=18, bold=True, color=BLUE_500, align=PP_ALIGN.CENTER)
add_page_number(s, 9, color=TEXT_GRAY)

output_path = os.path.join(os.path.dirname(__file__), "QuizMaster_Presentation.pptx")
prs.save(output_path)
print(f"PPTX saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
